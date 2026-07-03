import os
import csv

from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_text(filepath):

    extension = os.path.splitext(filepath)[1].lower()

    # ---------------- TXT ----------------
    if extension == ".txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # ---------------- PDF ----------------
    elif extension == ".pdf":
        text = ""
        reader = PdfReader(filepath)

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        return text

    # ---------------- DOCX ----------------
    elif extension == ".docx":
        doc = Document(filepath)

        text = ""

        for para in doc.paragraphs:
            text += para.text + "\n"

        return text

    # ---------------- XLSX ----------------
    elif extension == ".xlsx":
        wb = load_workbook(filepath)

        text = ""

        for sheet in wb.worksheets:
            text += sheet.title + "\n"

            for row in sheet.iter_rows(values_only=True):
                text += " ".join(
                    str(cell) if cell else ""
                    for cell in row
                ) + "\n"

        return text

    # ---------------- CSV ----------------
    elif extension == ".csv":
        text = ""

        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)

            for row in reader:
                text += " ".join(row) + "\n"

        return text

    # ---------------- PPTX ----------------
    elif extension == ".pptx":

        prs = Presentation(filepath)

        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    # ---------------- Unsupported ----------------
    else:
        return "Unsupported file type."